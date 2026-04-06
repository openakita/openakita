"""
OpenAkita 企业定制 Agent 解决方案 PPT
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

IMG_DIR = r"D:\coder\myagent\ppt_images"
OUTPUT = r"D:\coder\myagent\OpenAkita_企业定制Agent解决方案.pptx"

ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ORANGE_LIGHT = RGBColor(0xFF, 0xA9, 0x40)
ORANGE_DARK = RGBColor(0xE6, 0x7E, 0x00)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
DARK = RGBColor(0x1A, 0x1A, 0x2E)
DARK_MID = RGBColor(0x16, 0x21, 0x3E)
DARK_LIGHT = RGBColor(0x0F, 0x34, 0x60)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_TEXT = RGBColor(0x71, 0x80, 0x96)
LIGHT_GRAY = RGBColor(0xF7, 0xFA, 0xFC)
BLUE = RGBColor(0x3B, 0x82, 0xF6)
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)
RED = RGBColor(0xEF, 0x44, 0x44)
CARD_BG = RGBColor(0x22, 0x22, 0x42)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H


def img(name):
    return os.path.join(IMG_DIR, name)


def set_bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def rrect(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def oval(slide, l, t, w, h, fill):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.fill.background()
    return s


def tb(slide, l, t, w, h, text, sz=18, clr=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    box.text_frame.word_wrap = True
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = clr
    p.font.bold = bold
    p.font.name = 'Arial'
    p.alignment = align
    return box


def multi_tb(slide, l, t, w, h, lines, sz=14, clr=WHITE, bold=False, align=PP_ALIGN.LEFT, sp=1.3):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, txt in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = 'Arial'
        p.alignment = align
        p.space_after = Pt(sz * (sp - 1))
    return box


def rich_tb(slide, l, t, w, h, paras, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for pi, runs in enumerate(paras):
        p = tf.paragraphs[0] if pi == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(4)
        for text, sz, clr, bold in runs:
            r = p.add_run()
            r.text = text
            r.font.size = Pt(sz)
            r.font.color.rgb = clr
            r.font.bold = bold
            r.font.name = 'Arial'
    return box


def tag(slide, l, t, w, h, text, bg_clr, txt_clr):
    s = rrect(slide, l, t, w, h, bg_clr)
    p = s.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(11)
    p.font.color.rgb = txt_clr
    p.font.bold = True
    p.font.name = 'Arial'
    p.alignment = PP_ALIGN.CENTER
    return s


def top_bar(slide, color=ORANGE):
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), color)


def page_num(slide, num, total, color=GRAY_TEXT):
    tb(slide, Inches(12.0), Inches(7.1), Inches(1.2), Inches(0.3),
       f"{num:02d} / {total:02d}", 10, color, False, PP_ALIGN.RIGHT)


TOTAL = 13

# ============================================================
# 1 - 封面
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)

s.shapes.add_picture(img("cover_hero.png"), Inches(6.2), Inches(0.5), Inches(6.8), Inches(3.8))
s.shapes.add_picture(img("whitelabel.png"), Inches(5.5), Inches(3.8), Inches(7.5), Inches(3.5))

rect(s, Inches(0), Inches(0), Inches(7), SLIDE_H, DARK)

tb(s, Inches(0.8), Inches(0.6), Inches(3), Inches(0.4), "OpenAkita", 16, ORANGE, True)

rich_tb(s, Inches(0.8), Inches(1.6), Inches(5.5), Inches(2.0), [
    [("企业定制 Agent", 44, WHITE, True)],
    [("解决方案", 44, ORANGE, True)],
])

multi_tb(s, Inches(0.8), Inches(3.6), Inches(5.5), Inches(1.0), [
    "基于国产全栈自研开源多 Agent AI 助手",
    "为您打造专属的本地化智能助手系统",
], 18, RGBColor(0xA0, 0xAE, 0xC0))

tags_data = ["前端定制开发", "定向 LLM 服务", "算力接入", "功能深度定制", "私有化部署"]
for i, t in enumerate(tags_data):
    tag(s, Inches(0.8 + i * 2.1), Inches(5.2), Inches(1.9), Inches(0.42),
        t, RGBColor(0x2A, 0x2A, 0x4E), ORANGE_LIGHT)

tb(s, Inches(0.8), Inches(6.4), Inches(6), Inches(0.3),
   "openakita.ai  |  github.com/openakita  |  Apache 2.0 License",
   11, RGBColor(0x50, 0x50, 0x70))
page_num(s, 1, TOTAL, RGBColor(0x50, 0x50, 0x70))


# ============================================================
# 2 - 关于 OpenAkita
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(3), Inches(0.3), "ABOUT OPENAKITA", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "关于 OpenAkita — 不只是聊天，是帮你做事的 AI 团队", 28, DARK, True)

s.shapes.add_picture(img("customization_overview.png"),
                     Inches(7.0), Inches(1.4), Inches(5.8), Inches(3.2))

desc = ("OpenAkita 是一款国产全栈自研的开源多 Agent AI 助手。"
        "将复杂的 AI 技术封装为极简体验——无需技术基础，三分钟自助安装。"
        "通过多 Agent 协作架构，AI 可以像一个团队一样理解、拆解并执行任务。")
tb(s, Inches(0.8), Inches(1.5), Inches(5.8), Inches(1.5), desc, 14, RGBColor(0x4A, 0x55, 0x68))

stats = [("89+", "内置工具", ORANGE), ("148+", "可用技能", TEAL),
         ("30+", "LLM 服务商", BLUE), ("6", "IM 平台", PURPLE), ("3 端", "全覆盖", ORANGE)]
for i, (num, label, clr) in enumerate(stats):
    x = Inches(0.8 + i * 1.2)
    card = rrect(s, x, Inches(3.3), Inches(1.05), Inches(1.1), LIGHT_GRAY)
    tb(s, x + Inches(0.05), Inches(3.4), Inches(0.95), Inches(0.5), num, 24, clr, True, PP_ALIGN.CENTER)
    tb(s, x + Inches(0.05), Inches(3.9), Inches(0.95), Inches(0.3), label, 10, GRAY_TEXT, False, PP_ALIGN.CENTER)

highlights = [
    ("数据本地存储", "记忆、配置、对话全部存在你自己的电脑上，数据不出企业"),
    ("国内深度适配", "通义千问/DeepSeek/Kimi 等国产 LLM，飞书/企微/钉钉 IM"),
    ("完全开源可控", "Apache 2.0 许可证，代码完全公开，自由修改、商用无忧"),
]
for i, (title, desc) in enumerate(highlights):
    x = Inches(0.8 + i * 4.1)
    card = rrect(s, x, Inches(4.8), Inches(3.8), Inches(1.2), RGBColor(0xF0, 0xF7, 0xFF))
    rect(s, x, Inches(4.8), Inches(0.05), Inches(1.2), ORANGE)
    tb(s, x + Inches(0.2), Inches(4.9), Inches(3.4), Inches(0.35), title, 14, DARK, True)
    tb(s, x + Inches(0.2), Inches(5.3), Inches(3.4), Inches(0.5), desc, 11, GRAY_TEXT)

s.shapes.add_picture(img("security_privacy.png"),
                     Inches(7.0), Inches(4.6), Inches(5.8), Inches(2.5))

tb(s, Inches(0.8), Inches(6.3), Inches(10), Inches(0.3),
   "全栈自研  |  Apache 2.0 开源  |  GitHub 1.3K+ Stars  |  多 Agent 协作  |  组织编排  |  3 分钟安装",
   10, GRAY_TEXT)
page_num(s, 2, TOTAL)


# ============================================================
# 3 - 为什么选择 OpenAkita
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "WHY OPENAKITA", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "为什么选择 OpenAkita 作为定制基座?", 28, DARK, True)

advantages = [
    ("100%", "全栈自研\n深度可控", "从身份系统到追踪层七层架构\n全部自主研发，代码完全透明\n可按需修改任何模块", ORANGE),
    ("30+", "国产生态\n深度适配", "原生支持通义千问/DeepSeek/\nKimi 等国产 LLM，飞书/企微/\n钉钉等 IM 平台深度集成", TEAL),
    ("20+", "多 Agent\n协作架构", "业界领先的多 Agent 协作与\n组织编排能力，5 层委派深度\n20+ 预置专业角色", BLUE),
    ("零风险", "私有化部署\n安全可控", "数据完全本地存储，无需上云\n满足等保/信创/内网部署等\n合规需求", PURPLE),
]

for i, (stat, title, desc, clr) in enumerate(advantages):
    x = Inches(0.6 + i * 3.15)
    card = rrect(s, x, Inches(1.5), Inches(2.9), Inches(4.2), LIGHT_GRAY)
    rect(s, x, Inches(1.5), Inches(2.9), Inches(0.05), clr)
    tb(s, x + Inches(0.3), Inches(1.7), Inches(2.3), Inches(0.6), stat, 36, clr, True)
    multi_tb(s, x + Inches(0.3), Inches(2.3), Inches(2.3), Inches(0.8),
             title.split('\n'), 15, DARK, True, PP_ALIGN.LEFT, 1.2)
    multi_tb(s, x + Inches(0.3), Inches(3.2), Inches(2.3), Inches(2.0),
             desc.split('\n'), 11, GRAY_TEXT, False, PP_ALIGN.LEFT, 1.5)

rect(s, Inches(0.5), Inches(6.0), Inches(12.3), Inches(1.0), RGBColor(0xFF, 0xF8, 0xF0))
rich_tb(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.8), [
    [("核心优势: ", 13, ORANGE_DARK, True),
     ("OpenAkita 是目前市面上", 13, RGBColor(0x4A, 0x55, 0x68), False),
     ("唯一同时具备「全栈自研 + 多 Agent 协作 + 组织编排 + 自我进化 + 全端 GUI + 数据本地化」", 13, DARK, True),
     ("的开源 AI 助手框架", 13, RGBColor(0x4A, 0x55, 0x68), False)],
])
page_num(s, 3, TOTAL)


# ============================================================
# 4 - 四大定制服务
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "CUSTOMIZATION SERVICES", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "四大定制服务 — 打造您的专属 AI 助手", 28, WHITE, True)
tb(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.4),
   "从界面到内核，从算力到功能，全方位深度定制", 15, RGBColor(0xA0, 0xAE, 0xC0))

services = [
    ("01", "前端定制开发", "品牌 UI/UX 深度定制",
     ["品牌 Logo & 主题色替换", "定制桌面端/Web端/移动端", "专属交互界面设计",
      "多语言本地化支持", "定制仪表盘与管理后台"], ORANGE),
    ("02", "定向 LLM 服务", "接入您自己的大模型",
     ["对接客户自有 LLM/私有算力", "内置指定服务商模型", "智能路由与故障切换",
      "Token 用量与成本优化", "离线/内网 LLM 部署"], TEAL),
    ("03", "算力接入", "灵活对接计算资源",
     ["对接自有 GPU 集群", "主流云厂商算力对接", "本地 Ollama/vLLM 部署",
      "弹性调度与负载均衡", "资源预算与成本控制"], BLUE),
    ("04", "功能深度定制", "按需扩展核心能力",
     ["自定义 Agent 角色与技能", "行业专属工具链开发", "定制组织编排流程",
      "专属 IM 通道集成", "定制记忆系统与知识库"], PURPLE),
]

for i, (num, title, subtitle, items, clr) in enumerate(services):
    x = Inches(0.5 + i * 3.15)
    y = Inches(2.0)
    card = rrect(s, x, y, Inches(2.95), Inches(5.1), CARD_BG)
    tb(s, x + Inches(0.3), y + Inches(0.25), Inches(0.8), Inches(0.5), num, 26, clr, True)
    tb(s, x + Inches(0.3), y + Inches(0.85), Inches(2.3), Inches(0.4), title, 18, WHITE, True)
    tb(s, x + Inches(0.3), y + Inches(1.3), Inches(2.3), Inches(0.3), subtitle, 11, RGBColor(0xA0, 0xAE, 0xC0))
    rect(s, x + Inches(0.3), y + Inches(1.75), Inches(2.3), Inches(0.02), clr)
    for j, item in enumerate(items):
        tb(s, x + Inches(0.3), y + Inches(1.95 + j * 0.55), Inches(2.3), Inches(0.45),
           f"•  {item}", 11, RGBColor(0xC0, 0xCC, 0xDD))

page_num(s, 4, TOTAL, RGBColor(0x50, 0x50, 0x70))


# ============================================================
# 5 - 前端定制开发
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "FRONTEND CUSTOMIZATION", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "前端定制开发 — 您的品牌，您的界面", 28, DARK, True)

s.shapes.add_picture(img("frontend_custom.png"), Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.2))

features = [
    ("品牌形象定制", "替换 Logo、主题色、启动画面等全套品牌元素，打造完整企业 VI 体系"),
    ("桌面端深度定制", "基于 Tauri 2.x + React 18，Windows/macOS/Linux 三端原生体验"),
    ("移动端适配", "Capacitor 跨平台方案，Android/iOS 同时覆盖，支持原生推送"),
    ("Web 管理后台", "用户管理、Agent 监控、权限配置、使用统计等功能一应俱全"),
    ("定制化仪表盘", "Agent 状态可视化、Token 追踪、任务统计，支持自定义看板"),
]

for i, (title, desc) in enumerate(features):
    y = Inches(1.5 + i * 1.0)
    rrect(s, Inches(0.6), y, Inches(6.0), Inches(0.85), LIGHT_GRAY)
    rect(s, Inches(0.6), y, Inches(0.05), Inches(0.85), ORANGE)
    tb(s, Inches(0.9), y + Inches(0.05), Inches(5.4), Inches(0.3), title, 14, DARK, True)
    tb(s, Inches(0.9), y + Inches(0.4), Inches(5.4), Inches(0.4), desc, 11, GRAY_TEXT)

techs = ["Tauri 2.x", "React 18", "TypeScript", "Vite 6", "Capacitor", "11 功能面板"]
for i, t in enumerate(techs):
    col, row = i % 3, i // 3
    tag(s, Inches(7.2 + col * 1.9), Inches(4.8 + row * 0.55), Inches(1.7), Inches(0.4),
        t, RGBColor(0xFF, 0xF3, 0xE0), ORANGE_DARK)

s.shapes.add_picture(img("whitelabel.png"), Inches(7.0), Inches(5.5), Inches(5.8), Inches(1.8))
page_num(s, 5, TOTAL)


# ============================================================
# 6 - LLM & 算力接入
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, TEAL)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "LLM & COMPUTE", 11, TEAL, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "定向 LLM 服务 & 算力接入 — 用您自己的大脑", 28, DARK, True)

s.shapes.add_picture(img("llm_compute.png"), Inches(0.3), Inches(1.4), Inches(5.5), Inches(3.0))

llm_items = [
    ("客户自有 LLM", "直接对接私有部署大模型: ChatGLM、Baichuan、Qwen 等"),
    ("云端 API 服务", "30+ 国内外 LLM 服务商一键对接，智能故障切换"),
    ("本地模型部署", "Ollama / vLLM / llama.cpp，GPU 算力不出内网"),
    ("混合路由策略", "按任务类型/成本/速度自动选最优模型"),
]
for i, (title, desc) in enumerate(llm_items):
    y = Inches(1.5 + i * 0.95)
    tb(s, Inches(6.2), y, Inches(6.0), Inches(0.3), f"▸ {title}", 14, DARK_LIGHT, True)
    tb(s, Inches(6.4), y + Inches(0.3), Inches(5.8), Inches(0.5), desc, 11, GRAY_TEXT)

rect(s, Inches(0.3), Inches(4.6), Inches(12.7), Inches(0.03), RGBColor(0xE2, 0xE8, 0xF0))

compute_items = [
    ("自有 GPU 集群", "直接对接 NVIDIA A100/H100 等 GPU 集群"),
    ("公有云算力", "阿里云 PAI、华为云 ModelArts、腾讯云 TI 等"),
    ("边缘计算", "Jetson 等边缘设备部署，适用 IoT/工厂离线场景"),
    ("弹性伸缩", "根据负载自动调整算力，资源利用率最大化"),
]
for i, (title, desc) in enumerate(compute_items):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 6.3)
    y = Inches(4.9 + row * 1.15)
    rrect(s, x, y, Inches(6.0), Inches(1.0), RGBColor(0xF0, 0xFF, 0xF4))
    rect(s, x, y, Inches(0.05), Inches(1.0), TEAL)
    tb(s, x + Inches(0.2), y + Inches(0.08), Inches(5.6), Inches(0.3), title, 13, DARK, True)
    tb(s, x + Inches(0.2), y + Inches(0.42), Inches(5.6), Inches(0.4), desc, 11, GRAY_TEXT)

page_num(s, 6, TOTAL)


# ============================================================
# 7 - 功能深度定制
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s, PURPLE)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "FEATURE CUSTOMIZATION", 11, PURPLE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "功能深度定制 — 让 AI 精准匹配您的业务", 28, DARK, True)

s.shapes.add_picture(img("multi_agent.png"), Inches(0.3), Inches(1.4), Inches(6.2), Inches(3.2))

cats = [
    ("Agent 角色定制", ["定制专属 Agent 角色", "配置角色三元组 Title/Goal/Backstory",
                       "自定义人格与行为准则", "行业知识库注入"], ORANGE),
    ("工具链 & 技能定制", ["行业专属工具开发 (CRM/ERP等)", "MCP 协议工具扩展",
                         "SKILL.md 声明式技能定制", "技能市场私有部署"], BLUE),
    ("组织编排定制", ["模拟客户真实组织架构", "定制部门协作与审批链",
                    "黑板共享记忆三级系统", "自主运营与心跳监控"], TEAL),
    ("IM 通道 & 集成", ["企微/飞书/钉钉机器人定制", "Telegram Bot 私有化部署",
                       "自定义消息格式", "第三方 Webhook 集成"], PURPLE),
]

for i, (title, items, clr) in enumerate(cats):
    col, row = i % 2, i // 2
    x = Inches(6.8 + col * 3.15)
    y = Inches(1.4 + row * 2.85)
    card = rrect(s, x, y, Inches(3.0), Inches(2.7), LIGHT_GRAY)
    rect(s, x, y, Inches(3.0), Inches(0.04), clr)
    tb(s, x + Inches(0.2), y + Inches(0.12), Inches(2.6), Inches(0.35), title, 14, DARK, True)
    for j, item in enumerate(items):
        tb(s, x + Inches(0.2), y + Inches(0.55 + j * 0.45), Inches(2.6), Inches(0.4),
           f"•  {item}", 10, GRAY_TEXT)

page_num(s, 7, TOTAL)


# ============================================================
# 8 - 技术架构
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)
top_bar(s, BLUE)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "TECHNICAL ARCHITECTURE", 11, BLUE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "七层全栈技术架构 — 每一层都可定制", 28, WHITE, True)

s.shapes.add_picture(img("architecture.png"), Inches(0.3), Inches(1.4), Inches(5.8), Inches(3.8))

layers = [
    ("桌面应用", "Tauri 2.x / React 18 / TypeScript / 11 功能面板", ORANGE),
    ("身份层", "SOUL.md / AGENT.md / POLICIES.yaml / 8 种人格", PURPLE),
    ("核心层", "ReasoningEngine (ReAct) / Brain / ContextManager", BLUE),
    ("Agent 层", "Orchestrator / InstancePool / Factory / AgentOrg", TEAL),
    ("记忆层", "UnifiedStore (SQLite+向量) / Retrieval / Extractor", ORANGE_LIGHT),
    ("工具层", "Shell / File / Browser / Desktop / Web / MCP / Skills", RED),
    ("通道层", "CLI / Telegram / 飞书 / 企微 / 钉钉 / QQ / OneBot", PURPLE),
]

for i, (name, items, clr) in enumerate(layers):
    y = Inches(1.5 + i * 0.78)
    rrect(s, Inches(6.5), y, Inches(6.3), Inches(0.65), CARD_BG)
    rect(s, Inches(6.5), y, Inches(0.05), Inches(0.65), clr)
    tb(s, Inches(6.7), y + Inches(0.03), Inches(1.5), Inches(0.3), name, 12, clr, True)
    tb(s, Inches(6.7), y + Inches(0.3), Inches(5.9), Inches(0.3), items, 9, RGBColor(0xA0, 0xAE, 0xC0))

techs = [
    ("后端", "Python 3.11+ / FastAPI / asyncio"),
    ("桌面端", "Tauri 2.x (Rust) / React 18"),
    ("数据", "SQLite + 向量数据库"),
    ("协议", "MCP (Model Context Protocol)"),
]
for i, (k, v) in enumerate(techs):
    col = i % 2
    row = i // 2
    x = Inches(0.5 + col * 3.0)
    y = Inches(5.6 + row * 0.65)
    tb(s, x, y, Inches(1.0), Inches(0.3), k, 11, ORANGE_LIGHT, True)
    tb(s, x + Inches(1.0), y, Inches(2.0), Inches(0.3), v, 10, RGBColor(0xA0, 0xAE, 0xC0))

page_num(s, 8, TOTAL, RGBColor(0x50, 0x50, 0x70))


# ============================================================
# 9 - 行业场景
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "INDUSTRY SCENARIOS", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "行业场景定制 — AI 团队赋能千行百业", 28, DARK, True)

s.shapes.add_picture(img("industry_scenes.png"), Inches(0.3), Inches(1.4), Inches(6.0), Inches(3.3))

scenes = [
    ("跨境电商", "AI 选品中台，9 Agent 并行\n市场发现/合规审查/竞争分析", "效率提升 10 倍", ORANGE),
    ("智慧监测", "设备日志+运行数据 AI 分析\n提前 72h 预测故障", "7x24 无人值守", BLUE),
    ("电商运营", "订单/客服/报表/库存预警\n全流程 AI 驱动", "减少 50% 重复劳动", TEAL),
    ("法务合规", "合同审查/条款风险识别\n法规检索/报告生成", "审查效率 8 倍", PURPLE),
    ("教育培训", "学习计划/答疑/批改\n个性化教学策略", "记忆追踪学情", ORANGE),
    ("内容创作", "选题/写作/SEO 优化\n多平台自动分发", "全流程自动化", BLUE),
]

for i, (title, desc, result, clr) in enumerate(scenes):
    col, row = i % 2, i // 2
    x = Inches(6.5 + col * 3.3)
    y = Inches(1.4 + row * 1.85)
    card = rrect(s, x, y, Inches(3.1), Inches(1.7), LIGHT_GRAY)
    rect(s, x, y, Inches(3.1), Inches(0.04), clr)
    tb(s, x + Inches(0.2), y + Inches(0.1), Inches(2.7), Inches(0.3), title, 15, DARK, True)
    multi_tb(s, x + Inches(0.2), y + Inches(0.45), Inches(2.7), Inches(0.7),
             desc.split('\n'), 10, GRAY_TEXT, False, PP_ALIGN.LEFT, 1.3)
    tag(s, x + Inches(0.2), y + Inches(1.2), Inches(2.7), Inches(0.3), f"▸ {result}", clr, WHITE)

more = ["制造业 · 生产管理", "医疗 · 辅助诊疗", "财务 · 智能核算", "HR · 招聘管理"]
for i, t in enumerate(more):
    tag(s, Inches(0.5 + i * 3.15), Inches(6.8), Inches(2.9), Inches(0.4),
        t, RGBColor(0xF0, 0xF7, 0xFF), DARK_LIGHT)

page_num(s, 9, TOTAL)


# ============================================================
# 10 - 交付流程
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "DELIVERY PROCESS", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "项目交付流程 — 从需求到上线一站式服务", 28, WHITE, True)

s.shapes.add_picture(img("delivery_process.png"), Inches(0.3), Inches(1.3), Inches(12.7), Inches(3.0))

steps = [
    ("01", "需求调研", "深入了解业务\n场景与痛点", "1-2 周", ORANGE),
    ("02", "方案设计", "架构方案\n技术选型排期", "1 周", TEAL),
    ("03", "定制开发", "前端/后端/LLM\n并行开发", "4-8 周", BLUE),
    ("04", "联调测试", "全场景测试\n性能优化", "1-2 周", PURPLE),
    ("05", "部署上线", "私有化部署\n运维配置", "1 周", RED),
    ("06", "持续服务", "技术支持\n版本升级", "持续", RGBColor(0x10, 0xB9, 0x81)),
]

for i, (num, title, desc, dur, clr) in enumerate(steps):
    x = Inches(0.5 + i * 2.12)
    y = Inches(4.3)
    card = rrect(s, x, y, Inches(1.95), Inches(2.7), CARD_BG)
    circle = oval(s, x + Inches(0.6), y + Inches(0.2), Inches(0.7), Inches(0.7), clr)
    p = circle.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tb(s, x + Inches(0.1), y + Inches(1.0), Inches(1.75), Inches(0.35), title, 15, WHITE, True, PP_ALIGN.CENTER)
    multi_tb(s, x + Inches(0.1), y + Inches(1.4), Inches(1.75), Inches(0.6),
             desc.split('\n'), 10, RGBColor(0xA0, 0xAE, 0xC0), False, PP_ALIGN.CENTER, 1.3)
    tag(s, x + Inches(0.3), y + Inches(2.15), Inches(1.35), Inches(0.3), dur, clr, WHITE)

page_num(s, 10, TOTAL, RGBColor(0x50, 0x50, 0x70))


# ============================================================
# 11 - 合作模式
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "ENGAGEMENT MODELS", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "灵活的合作模式 — 按需选择最适合您的方案", 28, DARK, True)

s.shapes.add_picture(img("service_tiers.png"), Inches(0.3), Inches(1.3), Inches(12.7), Inches(2.2))

models = [
    ("轻量定制", "品牌换肤 + 基础配置",
     ["Logo 与品牌色替换", "指定 LLM 服务商", "基础 Agent 配置",
      "单一 IM 通道集成", "3 个月技术支持"],
     "适合: 快速上线品牌 Agent", ORANGE),
    ("标准定制", "功能定制 + 多通道",
     ["完整品牌 UI 定制", "多 LLM + 混合路由", "5-10 个 Agent 角色",
      "多 IM 通道集成", "6 个月技术支持"],
     "适合: 中型企业业务场景", TEAL),
    ("深度定制", "全栈深度开发",
     ["全套前端深度定制", "私有 LLM + 算力集群", "行业工具链开发",
      "组织编排定制", "12 个月支持 + SLA"],
     "适合: 大型企业私有化部署", BLUE),
    ("战略合作", "联合开发 + 长期运营",
     ["全部深度定制能力", "联合产品规划", "专属开发团队",
      "持续迭代与共创", "无限期技术支持"],
     "适合: 战略级 AI 转型伙伴", PURPLE),
]

for i, (title, subtitle, items, note, clr) in enumerate(models):
    x = Inches(0.5 + i * 3.15)
    y = Inches(3.6)
    card = rrect(s, x, y, Inches(2.95), Inches(3.7), LIGHT_GRAY)
    rect(s, x, y, Inches(2.95), Inches(0.05), clr)
    tb(s, x + Inches(0.2), y + Inches(0.15), Inches(2.5), Inches(0.4), title, 18, DARK, True, PP_ALIGN.CENTER)
    tb(s, x + Inches(0.2), y + Inches(0.55), Inches(2.5), Inches(0.3), subtitle, 11, GRAY_TEXT, False, PP_ALIGN.CENTER)
    rect(s, x + Inches(0.2), y + Inches(0.95), Inches(2.5), Inches(0.01), RGBColor(0xE2, 0xE8, 0xF0))
    for j, item in enumerate(items):
        tb(s, x + Inches(0.2), y + Inches(1.1 + j * 0.4), Inches(2.5), Inches(0.35),
           f"  {item}", 10, RGBColor(0x4A, 0x55, 0x68))
    tag(s, x + Inches(0.15), y + Inches(3.2), Inches(2.65), Inches(0.35), note, clr, WHITE)

page_num(s, 11, TOTAL)


# ============================================================
# 12 - 定制案例
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, WHITE)
top_bar(s)

tb(s, Inches(0.8), Inches(0.35), Inches(5), Inches(0.3), "CASE STUDIES", 11, ORANGE, True)
tb(s, Inches(0.8), Inches(0.7), Inches(10), Inches(0.6),
   "定制场景参考 — 真实业务落地", 28, DARK, True)

s.shapes.add_picture(img("case_studies.png"), Inches(0.3), Inches(1.3), Inches(4.5), Inches(2.5))

cases = [
    ("跨境电商 AI 选品中台", "大型跨境电商",
     "9 Agent 并行: 市场发现、合规审查、竞争分析、盈利测算，飞书实时推送选品报告",
     ["组织编排", "飞书集成", "自定义工具"], "效率提升 10 倍", ORANGE),
    ("智慧监测 AI 运维助手", "工业制造企业",
     "实时采集设备日志，AI 趋势分析+故障预测，自动生成报告，企微即时告警",
     ["私有化部署", "企微集成", "IoT 工具"], "提前 72h 预警", BLUE),
    ("电商运营全流程", "中小电商企业",
     "订单处理、智能客服、销售报表、库存预警，钉钉实时推送，一人管理全店",
     ["Agent 定制", "钉钉集成", "定制报表"], "减少 50% 重复劳动", TEAL),
    ("法务 AI 合同审查", "律师事务所",
     "合同逐条扫描风险，对照法规出修改建议，内置法律知识库，数据零泄露",
     ["知识库定制", "私有 LLM", "本地部署"], "审查效率 8 倍", PURPLE),
]

for i, (title, client, desc, tags_list, result, clr) in enumerate(cases):
    col, row = i % 2, i // 2
    x = Inches(5.0 + col * 4.1)
    y = Inches(1.3 + row * 2.95)
    card = rrect(s, x, y, Inches(3.9), Inches(2.8), LIGHT_GRAY)
    rect(s, x, y, Inches(3.9), Inches(0.04), clr)
    tb(s, x + Inches(0.2), y + Inches(0.1), Inches(2.5), Inches(0.3), title, 14, DARK, True)
    tb(s, x + Inches(2.8), y + Inches(0.15), Inches(0.9), Inches(0.2), client, 9, GRAY_TEXT, False, PP_ALIGN.RIGHT)
    tb(s, x + Inches(0.2), y + Inches(0.45), Inches(3.5), Inches(0.8), desc, 10, RGBColor(0x4A, 0x55, 0x68))
    for j, t in enumerate(tags_list):
        tag(s, x + Inches(0.2 + j * 1.2), y + Inches(1.35), Inches(1.1), Inches(0.25),
            t, RGBColor(0xFF, 0xF3, 0xE0), ORANGE_DARK)
    tag(s, x + Inches(0.2), y + Inches(1.75), Inches(3.5), Inches(0.3),
        f"▸ 成果: {result}", clr, WHITE)

guarantees = ["源码交付", "技术培训", "SLA 保障", "版本升级", "安全审计", "运维手册"]
for i, g in enumerate(guarantees):
    tag(s, Inches(0.5 + i * 2.12), Inches(7.0), Inches(1.95), Inches(0.35),
        g, RGBColor(0xF0, 0xF7, 0xFF), DARK_LIGHT)

page_num(s, 12, TOTAL)


# ============================================================
# 13 - 联系我们
# ============================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(s, DARK)

s.shapes.add_picture(img("contact_us.png"), Inches(6.5), Inches(0.5), Inches(6.5), Inches(4.0))

rect(s, Inches(0), Inches(0), Inches(7), SLIDE_H, DARK)

tb(s, Inches(0.8), Inches(0.6), Inches(3), Inches(0.3), "OpenAkita", 16, ORANGE, True)

rich_tb(s, Inches(0.8), Inches(1.5), Inches(5.5), Inches(1.5), [
    [("让我们一起打造", 36, WHITE, True)],
    [("您的专属 ", 36, WHITE, True), ("AI 团队", 36, ORANGE, True)],
])

contacts = [
    ("官网", "openakita.ai"),
    ("GitHub", "github.com/openakita"),
    ("邮箱", "contact@openakita.ai"),
    ("X/Twitter", "@openakita"),
]
for i, (label, value) in enumerate(contacts):
    y = Inches(3.5 + i * 0.55)
    rich_tb(s, Inches(0.8), y, Inches(5), Inches(0.4), [
        [(f"{label}: ", 14, RGBColor(0xA0, 0xAE, 0xC0), False),
         (value, 14, WHITE, True)],
    ])

final_tags = ["全栈自研", "国产开源", "多 Agent", "组织编排", "私有化部署", "安全可控"]
for i, t in enumerate(final_tags):
    tag(s, Inches(0.8 + i * 1.7), Inches(5.8), Inches(1.5), Inches(0.38),
        t, RGBColor(0x2A, 0x2A, 0x4E), ORANGE_LIGHT)

tb(s, Inches(0.8), Inches(6.5), Inches(10), Inches(0.4),
   "OpenAkita — 不只是聊天，是帮你做事的 AI 团队  |  Apache 2.0 License",
   12, RGBColor(0x50, 0x50, 0x70))

s.shapes.add_picture(img("cover_hero.png"), Inches(7.0), Inches(4.5), Inches(5.5), Inches(2.8))

page_num(s, 13, TOTAL, RGBColor(0x50, 0x50, 0x70))


# ============================================================
prs.save(OUTPUT)
print(f"PPT saved: {OUTPUT}")
print(f"Total slides: {len(prs.slides)}")
