"""
媒体处理器

处理各种媒体内容:
- 图片理解 (Vision)
- 文件内容提取 (PDF, Office, etc.)
"""

import logging
import mimetypes
from pathlib import Path
from typing import Any

from ..types import MediaFile, MediaStatus

logger = logging.getLogger(__name__)


class MediaHandler:
    """
    媒体处理器

    提供统一的媒体处理接口，支持:
    - 图片描述/理解
    - 文档内容提取
    """

    def __init__(
        self,
        brain: Any | None = None,
        enable_ocr: bool = True,
    ):
        """
        Args:
            brain: Brain 实例（用于图片理解）
            enable_ocr: 是否启用 OCR
        """
        self.brain = brain
        self.enable_ocr = enable_ocr
        self._ocr = None

    async def process(self, media: MediaFile) -> MediaFile:
        """
        处理媒体文件

        根据类型自动选择处理方式

        Args:
            media: 媒体文件

        Returns:
            处理后的媒体文件（带有 transcription/description/extracted_text）
        """
        if not media.local_path:
            logger.warning(f"Media {media.id} has no local path, skipping processing")
            return media

        try:
            if media.is_image:
                await self.describe_image(media)
            elif media.is_document:
                await self.extract_text(media)

            media.status = MediaStatus.PROCESSED

        except Exception as e:
            logger.error(f"Failed to process media {media.id}: {e}")

        return media

    async def describe_image(self, media: MediaFile) -> str:
        """
        图片理解/描述

        使用 Claude Vision 或其他多模态模型

        Args:
            media: 图片文件

        Returns:
            图片描述
        """
        if not media.local_path:
            raise ValueError("Media has no local path")

        logger.info(f"Describing image: {media.filename}")

        try:
            if self.brain:
                # 使用 Claude Vision
                description = await self._describe_with_vision(media.local_path)
            else:
                # 回退：使用 OCR
                description = await self._ocr_image(media.local_path)
        except Exception as e:
            logger.warning(f"Image description failed: {e}")
            description = f"[图片: {media.filename}]"

        media.description = description
        return description

    async def _describe_with_vision(self, image_path: str) -> str:
        """使用 Claude Vision 描述图片"""
        import base64

        # 读取图片并转 base64
        with open(image_path, "rb") as f:
            image_data = base64.standard_b64encode(f.read()).decode()

        # 确定 MIME 类型
        mime_type = mimetypes.guess_type(image_path)[0] or "image/jpeg"

        # 调用 Claude
        response = await self.brain.client.messages.create(
            model=self.brain.model,
            max_tokens=500,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {
                                "type": "base64",
                                "media_type": mime_type,
                                "data": image_data,
                            },
                        },
                        {
                            "type": "text",
                            "text": "请简要描述这张图片的内容，用中文回答。",
                        },
                    ],
                }
            ],
        )

        return response.content[0].text

    async def _ocr_image(self, image_path: str) -> str:
        """使用 OCR 提取图片文字"""
        if not self.enable_ocr:
            return ""

        try:
            # 尝试使用 pytesseract
            import pytesseract
            from PIL import Image

            image = Image.open(image_path)
            text = pytesseract.image_to_string(image, lang="chi_sim+eng")

            return text.strip() if text.strip() else "[图片无可识别文字]"

        except ImportError:
            from openakita.tools._import_helper import import_or_hint

            hint = import_or_hint("pytesseract")
            logger.warning(f"OCR 不可用: {hint}")
            return ""
        except Exception as e:
            logger.warning(f"OCR failed: {e}")
            return ""

    async def extract_text(self, media: MediaFile) -> str:
        """
        提取文件内容

        支持:
        - PDF
        - Office 文档 (docx, xlsx, pptx)
        - 文本文件

        Args:
            media: 文件

        Returns:
            提取的文本
        """
        if not media.local_path:
            raise ValueError("Media has no local path")

        logger.info(f"Extracting text from: {media.filename}")

        path = Path(media.local_path)
        extension = path.suffix.lower()

        try:
            if extension == ".pdf":
                text = await self._extract_pdf(path)
            elif extension == ".docx":
                text = await self._extract_docx(path)
            elif extension == ".xlsx":
                text = await self._extract_xlsx(path)
            elif extension == ".pptx":
                text = await self._extract_pptx(path)
            elif extension in (".txt", ".md", ".json", ".py", ".js", ".html", ".css"):
                text = path.read_bytes().decode("utf-8", errors="ignore")
            else:
                text = f"[文件: {media.filename}，不支持内容提取]"
        except Exception as e:
            logger.warning(f"Text extraction failed: {e}")
            text = f"[文件: {media.filename}，提取失败]"

        media.extracted_text = text
        return text

    async def _extract_pdf(self, path: Path) -> str:
        """提取 PDF 内容"""
        try:
            import fitz  # PyMuPDF

            doc = fitz.open(str(path))
            text_parts = []

            for page in doc:
                text_parts.append(page.get_text())

            doc.close()
            return "\n".join(text_parts)

        except ImportError:
            # 回退方案
            try:
                import pypdf

                reader = pypdf.PdfReader(str(path))
                text_parts = []

                for page in reader.pages:
                    text_parts.append(page.extract_text())

                return "\n".join(text_parts)

            except ImportError:
                from openakita.tools._import_helper import import_or_hint

                hint = import_or_hint("fitz")
                raise ImportError(f"PDF 提取不可用: {hint}")

    async def _extract_docx(self, path: Path) -> str:
        """提取 Word 文档内容"""
        try:
            from docx import Document

            doc = Document(str(path))
            text_parts = []

            for para in doc.paragraphs:
                text_parts.append(para.text)

            return "\n".join(text_parts)

        except ImportError:
            from openakita.tools._import_helper import import_or_hint

            hint = import_or_hint("docx")
            raise ImportError(f"DOCX 提取不可用: {hint}")

    async def _extract_xlsx(self, path: Path) -> str:
        """提取 Excel 内容"""
        try:
            import openpyxl

            wb = openpyxl.load_workbook(str(path), read_only=True)
            text_parts = []

            for sheet in wb.worksheets:
                text_parts.append(f"## Sheet: {sheet.title}")

                for row in sheet.iter_rows(max_row=100):  # 限制行数
                    cells = [str(cell.value) if cell.value else "" for cell in row]
                    text_parts.append(" | ".join(cells))

            wb.close()
            return "\n".join(text_parts)

        except ImportError:
            from openakita.tools._import_helper import import_or_hint

            hint = import_or_hint("openpyxl")
            raise ImportError(f"XLSX 提取不可用: {hint}")

    async def _extract_pptx(self, path: Path) -> str:
        """提取 PowerPoint 内容"""
        try:
            from pptx import Presentation

            prs = Presentation(str(path))
            text_parts = []

            for i, slide in enumerate(prs.slides):
                text_parts.append(f"## Slide {i + 1}")

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)

            return "\n".join(text_parts)

        except ImportError:
            from openakita.tools._import_helper import import_or_hint

            hint = import_or_hint("pptx")
            raise ImportError(f"PPTX 提取不可用: {hint}")
